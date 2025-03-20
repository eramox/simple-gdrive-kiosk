import os
import re
import yaml
import subprocess
import tempfile
from google.oauth2.credentials import Credentials
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
from io import BytesIO
import pptx
from PIL import Image

# Create output directory
output_dir = "slideshow_images"
os.makedirs(output_dir, exist_ok=True)

# Function to extract Google Drive file ID from URL
def extract_file_id(url):
    match = re.search(r'/d/([a-zA-Z0-9_-]+)', url)
    if match:
        return match.group(1)
    return None

# Function to download file from Google Drive
def download_file(file_id, output_file):
    # This requires OAuth2 setup - simplified here
    # In real implementation, you would need to set up proper authentication
    print(f"Downloading file with ID: {file_id}")
    # For demonstration, using wget with a direct link
    url = f"https://docs.google.com/document/d/{file_id}/export?format=docx"
    subprocess.run(["wget", "-O", output_file, url])

# Function to extract URL from docx file
def extract_url_from_docx(docx_file):
    # Convert docx to text and search for Google Slides URL
    temp_file = os.path.splitext('/home/user/somefile.txt')[0] + ''
    with open() as temp:
    	convert_cmd = ["libreoffice", "--headless", "--convert-to", "txt", 
                       "--outdir", os.path.dirname(temp.name), docx_file]
        print(f"{convert_cmd=}")
        subprocess.run(convert_cmd)
        txt_file = os.path.splitext(docx_file)[0] + '.txt'
        with open(txt_file, 'r') as f:
            content = f.read()
        
        # Look for Google Slides URL
        match = re.search(r'https://docs\.google\.com/presentation/d/([a-zA-Z0-9_-]+)', content)
        if match:
            return f"https://docs.google.com/presentation/d/{match.group(1)}"
    return None

# Function to download slides and extract durations
def download_and_process_slides(slides_url):
    slide_id = extract_file_id(slides_url)
    pptx_file = "presentation.pptx"
    
    # Download the presentation
    download_url = f"https://docs.google.com/presentation/d/{slide_id}/export/pptx"
    subprocess.run(["wget", "-O", pptx_file, download_url])
    
    # Process the presentation
    presentation = pptx.Presentation(pptx_file)
    slide_info = []
    
    for i, slide in enumerate(presentation.slides):
        # Extract slide notes to find duration info
        notes = slide.notes_slide.notes_text_frame.text if slide.has_notes_slide else ""
        duration = 5  # Default duration
        
        # Look for t=X format in notes
        match = re.search(r't=(\d+)', notes)
        if match:
            duration = int(match.group(1))
        
        # Save slide as image
        image_path = f"{output_dir}/slide_{i+1:03d}.jpg"
        
        # Export slide as PNG first (using LibreOffice)
        temp_dir = tempfile.mkdtemp()
        with open(f"{temp_dir}/slide_{i+1:03d}.odp", "wb") as odp_file:
            # Create a single-slide presentation
            # This is a simplified approach - in practice, you'd extract each slide
            subprocess.run(["libreoffice", "--headless", "--convert-to", "png", 
                           "--outdir", output_dir, f"{temp_dir}/slide_{i+1:03d}.odp"])
        
        # Convert PNG to JPG
        png_file = f"{output_dir}/slide_{i+1:03d}.png"
        if os.path.exists(png_file):
            Image.open(png_file).convert('RGB').save(image_path, 'JPEG')
            os.remove(png_file)  # Clean up PNG
        
        slide_info.append({
            "file": image_path,
            "duration": duration
        })
    
    # Create YAML file with slide info
    with open(f"{output_dir}/slideshow_config.yaml", "w") as yaml_file:
        yaml.dump({"slides": slide_info}, yaml_file, default_flow_style=False)
    
    print(f"Slideshow processed. {len(slide_info)} slides saved to {output_dir}")

# Main execution
doc_url = "https://docs.google.com/document/d/11hGse_OyqGQiy3qAaDMSHnrFX0_YyG2o/edit?usp=drive_link&ouid=104564737776631583414&rtpof=true&sd=true"
doc_id = extract_file_id(doc_url)
doc_file = "document.docx"

# Step 1: Download the document
download_file(doc_id, doc_file)

# Step 2: Extract the slides URL
slides_url = extract_url_from_docx(doc_file)
if slides_url:
    print(f"Found slides URL: {slides_url}")
    
    # Step 3 & 4: Download and process the slides
    download_and_process_slides(slides_url)
else:
    print("No slides URL found in the document")