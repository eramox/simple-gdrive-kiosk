import pptx
import lxml
import xml.dom.minidom
from pptx.oxml import parse_xml

ppt = pptx.Presentation('flprez.ppt')
slides = ppt.slides
print(f"{slides} {len(slides)}({type(slides)})")

# parser = lxml.etree.XMLParser(remove_blank_text=True)

for idx, slide in enumerate(slides):
    print(f"Slide {idx}= {slide}")
    elms = slide.element

    for elm in elms:
        print(f"elm {elm} {type(elm)}")

    str_blob = slide.part.blob
    ct_slide = parse_xml(str_blob)
    print(f"obj {ct_slide} {type(ct_slide)}")
    timing = ct_slide.get_or_add_childTnLst()
    print(f"timing {timing} {type(timing)}")
    # print(f"obj {ct_slide.timing}")
    # pretty_xml_as_string = root_element.toprettyxml()
    # print(pretty_xml_as_string)
    print("hello")
    # # CT_SlideTiming

    # # Search for text and read it
    # for sid, shape in enumerate(slide.shapes):
    #     print(f"shape {sid}: {shape}")
    #     if not shape.has_text_frame:
    #         continue
    #     text = shape.text_frame.text
    #     print(f"text found: {text}")

# self._rPr.get_or_add_latin()
# -> set <a:latin typeface="Arial"/>